import os
import uuid
import shutil
from fastapi import UploadFile
from typing import Dict, Any, List, Optional, Tuple
import asyncio
from datetime import datetime
import json
import logging  # Add logging import
import numpy as np  # Add numpy for image processing
import tiktoken  # Add tiktoken for token counting

# File processing libraries
import pytesseract
import fitz  # PyMuPDF
from pdf2image import convert_from_path
from PIL import Image
import docx
import openpyxl
from pptx import Presentation
from moviepy.video.io.VideoFileClip import VideoFileClip
from moviepy.audio.io.AudioFileClip import AudioFileClip
import whisper
from pytube import YouTube
import cv2  # Add OpenCV for frame extraction

# RAG specific imports
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_core.documents import Document as LangchainDocument

# Local imports
from rag.models import Document, ProcessingStatus, FileType, DocumentResponse
from rag.vector_store import add_chunks_to_vectorstore
from rag.database import save_document, update_document_status
from rag.repository import get_repository_by_id
from rag.video_processor import VideoProcessor  # Import our new VideoProcessor
from rag.config import (
    CHUNK_SIZE,
    CHUNK_OVERLAP,
    MAX_CHUNKS_PER_BATCH,
    MAX_TOKENS_PER_REQUEST
)

# Global variables
UPLOAD_FOLDER = os.getenv("UPLOAD_FOLDER", "./uploads")
PROCESSED_FOLDER = os.getenv("PROCESSED_FOLDER", "./processed")
FRAME_INTERVAL = 5  # Extract a frame every 5 seconds

# Initialize the Whisper model for audio transcription
whisper_model = whisper.load_model("base")

# Initialize the tokenizer for counting tokens
tokenizer = tiktoken.get_encoding("cl100k_base")  # OpenAI's embedding model tokenizer

# Configure logging
logger = logging.getLogger(__name__)
logger.setLevel(logging.DEBUG)

# Create console handler
console_handler = logging.StreamHandler()
console_handler.setLevel(logging.DEBUG)

# Create formatter
formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')
console_handler.setFormatter(formatter)

# Add handler to logger
logger.addHandler(console_handler)

async def process_document(
    file: UploadFile, 
    document_id: str, 
    repository_id: str,
    user_id: str,
    metadata: Dict[Any, Any]
) -> Dict[str, Any]:
    """
    Process an uploaded document, extract text, create chunks, and store in the vector database.
    
    Args:
        file: Uploaded file
        document_id: ID of the document
        repository_id: ID of the repository to add the document to
        user_id: ID of the user uploading the document
        metadata: Document metadata
        
    Returns:
        Dictionary with document information
    """
    # Verify repository exists
    repository = await get_repository_by_id(repository_id)
    if not repository:
        raise ValueError(f"Repository with ID {repository_id} does not exist")
    
    # Verify user has access to repository
    if repository.user_id != user_id:
        raise ValueError(f"User {user_id} does not have access to repository {repository_id}")
    
    # Save the uploaded file
    file_path = os.path.join(UPLOAD_FOLDER, f"{document_id}_{file.filename}")
    
    with open(file_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    # Create initial document record
    file_size = os.path.getsize(file_path)
    file_extension = file.filename.split(".")[-1].lower()
    
    file_type = determine_file_type(file_extension)
    
    document = Document(
        id=document_id,
        filename=file.filename,
        file_type=file_type,
        file_size=file_size,
        repository_id=repository_id,
        user_id=user_id,
        metadata=metadata,
        status=ProcessingStatus.PROCESSING,
        created_at=datetime.now(),
        updated_at=datetime.now()
    )
    
    # Save initial document to database
    await save_document(document)
    
    # Process the document asynchronously
    asyncio.create_task(
        process_document_async(document_id, file_path, file_type, repository_id, user_id)
    )
    
    return {
        "document_id": document_id,
        "filename": file.filename,
        "file_type": file_type,
        "repository_id": repository_id,
        "status": ProcessingStatus.PROCESSING.value,
        "message": "Document upload successful. Processing started."
    }

async def process_document_async(document_id: str, file_path: str, file_type: FileType, repository_id: str, user_id: str):
    """
    Asynchronously process a document and update its status.
    """
    try:
        # Extract text based on file type
        logger.info(f"Starting text extraction for document {document_id} ({file_type})")
        await update_document_status(document_id, ProcessingStatus.PROCESSING, "Extracting text content")
        
        text = await extract_text_from_file(file_path, file_type, user_id)
        
        # Create chunks
        logger.info(f"Creating chunks for document {document_id}")
        await update_document_status(document_id, ProcessingStatus.PROCESSING, "Creating text chunks")
        
        chunks = create_chunks(text, document_id, repository_id, user_id)
        
        # Calculate total tokens
        total_tokens = count_tokens_in_chunks(chunks)
        logger.info(f"Document {document_id} has {len(chunks)} chunks with approximately {total_tokens} tokens")
        
        # Create embeddings and store in vector database
        logger.info(f"Creating embeddings for document {document_id}")
        await update_document_status(document_id, ProcessingStatus.PROCESSING, "Creating embeddings")
        
        # Process chunks in batches to avoid token limits
        await process_chunks_in_batches(chunks, document_id)
        
        # Update document status to completed
        logger.info(f"Processing completed for document {document_id}")
        await update_document_status(document_id, ProcessingStatus.COMPLETED, "Processing completed successfully")
        
        # Move file to processed folder
        processed_path = os.path.join(PROCESSED_FOLDER, os.path.basename(file_path))
        shutil.move(file_path, processed_path)
        
    except Exception as e:
        logger.exception(f"Error processing document {document_id}: {str(e)}")
        await update_document_status(document_id, ProcessingStatus.FAILED, f"Processing failed: {str(e)}")

def determine_file_type(file_extension: str) -> FileType:
    """Determine the file type based on the file extension."""
    extension_map = {
        "pdf": FileType.PDF,
        "doc": FileType.DOC,
        "docx": FileType.DOCX,
        "xls": FileType.XLSX,
        "xlsx": FileType.XLSX,
        "ppt": FileType.PPT,
        "pptx": FileType.PPTX,
        "txt": FileType.TXT,
        "mp4": FileType.VIDEO,
        "mov": FileType.VIDEO,
        "avi": FileType.VIDEO,
        "mp3": FileType.AUDIO,
        "wav": FileType.AUDIO,
        "jpg": FileType.IMAGE,
        "jpeg": FileType.IMAGE,
        "png": FileType.IMAGE,
    }
    
    return extension_map.get(file_extension.lower(), FileType.UNKNOWN)

async def extract_text_from_file(file_path: str, file_type: FileType, user_id: str = None) -> str:
    """
    Extract text from different file types.
    
    Args:
        file_path: Path to the file
        file_type: Type of the file
        user_id: Optional user ID for text preprocessing
    """
    extracted_text = ""
    
    if file_type == FileType.PDF:
        try:
            extracted_text = await process_large_pdf_in_batches(file_path)
        except Exception as e:
            logger.exception(f"Error in batch PDF processing: {str(e)}")
            logger.info("Falling back to pdf2image method")
            # Use the pdf2image method as fallback
            try:
                text = ""
                # Convert PDF to images with lower resolution for speed
                images = convert_from_path(file_path, dpi=150)
                
                # Extract text from each image using OCR
                for i, image in enumerate(images):
                    text += pytesseract.image_to_string(image)
                    text += f"\nPage {i+1}\n"
                    # Free memory
                    del image
                
                logger.info("Fallback PDF extraction completed")
                extracted_text = text
            except Exception as e2:
                logger.exception(f"Fallback PDF extraction also failed: {str(e2)}")
                extracted_text = f"Error extracting text from PDF: {str(e)} / Fallback error: {str(e2)}"
    elif file_type in [FileType.DOC, FileType.DOCX]:
        extracted_text = extract_text_from_docx(file_path)
    elif file_type == FileType.XLSX:
        extracted_text = extract_text_from_xlsx(file_path)
    elif file_type in [FileType.PPT, FileType.PPTX]:
        extracted_text = extract_text_from_pptx(file_path)
    elif file_type == FileType.TXT:
        extracted_text = extract_text_from_txt(file_path)
    elif file_type == FileType.VIDEO:
        extracted_text = await extract_text_from_video(file_path)
    elif file_type == FileType.AUDIO:
        extracted_text = extract_text_from_audio(file_path)
    elif file_type == FileType.IMAGE:
        extracted_text = extract_text_from_image(file_path)
    else:
        extracted_text = "Unsupported file type"
    
    # Preprocess the extracted text if user_id is provided
    if user_id and extracted_text and not extracted_text.startswith("Error"):
        try:
            from rag.text_preprocessor import preprocess_document_text
            extracted_text = await preprocess_document_text(extracted_text, user_id)
        except Exception as e:
            logger.warning(f"Text preprocessing failed for user {user_id}: {str(e)}")
            # Continue with original text if preprocessing fails
    
    return extracted_text

async def process_large_pdf_in_batches(file_path: str, batch_size: int = 20) -> str:
    """
    Process a very large PDF in batches to prevent memory issues.
    
    Args:
        file_path: Path to the PDF file
        batch_size: Number of pages to process in each batch
        
    Returns:
        Extracted text from the PDF
    """
    logger.info(f"Starting batch processing for large PDF: {file_path}")
    
    try:
        # Open the PDF to get page count
        pdf_document = fitz.open(file_path)
        total_pages = pdf_document.page_count
        pdf_document.close()  # Close it immediately to free memory
        
        logger.info(f"Large PDF has {total_pages} pages, processing in batches of {batch_size}")
        
        # Calculate number of batches
        num_batches = (total_pages + batch_size - 1) // batch_size  # Ceiling division
        all_text = ""
        
        # Process PDF in batches
        for batch in range(num_batches):
            start_page = batch * batch_size
            end_page = min((batch + 1) * batch_size - 1, total_pages - 1)
            
            logger.info(f"Processing batch {batch+1}/{num_batches} (pages {start_page+1}-{end_page+1})")
            
            # Extract text from this batch of pages
            batch_text = await extract_pdf_page_range(file_path, start_page, end_page)
            all_text += batch_text
            
            # Update status (if needed)
            # await update_document_status(document_id, ProcessingStatus.PROCESSING, 
            #                            f"Processed {end_page+1} of {total_pages} pages")
            
            # Sleep briefly to allow other tasks to run
            await asyncio.sleep(0.1)
        
        logger.info(f"Completed batch processing of {total_pages} pages")
        return all_text
        
    except Exception as e:
        logger.exception(f"Error in batch PDF processing: {str(e)}")
        return f"Error processing PDF in batches: {str(e)}"

async def extract_pdf_page_range(file_path: str, start_page: int, end_page: int) -> str:
    """
    Extract text from a specific range of pages in a PDF.
    
    Args:
        file_path: Path to the PDF file
        start_page: First page to extract (0-indexed)
        end_page: Last page to extract (0-indexed)
        
    Returns:
        Extracted text from the specified range of pages
    """
    try:
        # Open the PDF
        pdf_document = fitz.open(file_path)
        
        # Validate page range
        if start_page < 0:
            start_page = 0
        if end_page >= pdf_document.page_count:
            end_page = pdf_document.page_count - 1
        
        # Extract text from pages
        text = ""
        pages_with_text = 0
        pages_needing_ocr = []
        
        # First try direct text extraction
        for page_num in range(start_page, end_page + 1):
            page = pdf_document[page_num]
            page_text = page.get_text()
            
            # If page has meaningful text content
            if len(page_text.strip()) > 50:
                text += f"\n--- Page {page_num + 1} ---\n"
                text += page_text
                pages_with_text += 1
            else:
                # This page might be scanned/image-based - add to OCR list
                pages_needing_ocr.append(page_num)
        
        # Process pages needing OCR if they're not too many
        pages_in_batch = end_page - start_page + 1
        if len(pages_needing_ocr) > 0 and (len(pages_needing_ocr) < pages_in_batch * 0.7):
            logger.info(f"Processing {len(pages_needing_ocr)} pages with OCR in batch {start_page+1}-{end_page+1}")
            
            for page_num in pages_needing_ocr:
                try:
                    page = pdf_document[page_num]
                    # Render page to image at a lower DPI for speed
                    pix = page.get_pixmap(matrix=fitz.Matrix(200/72, 200/72))  # 200 DPI
                    
                    # Convert to PIL Image
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    
                    # Perform OCR
                    ocr_text = pytesseract.image_to_string(img)
                    
                    if ocr_text.strip():
                        text += f"\n--- Page {page_num + 1} (OCR) ---\n"
                        text += ocr_text
                except Exception as e:
                    logger.warning(f"Error OCR processing page {page_num}: {str(e)}")
        
        # Close the PDF
        pdf_document.close()
        
        return text
        
    except Exception as e:
        logger.exception(f"Error extracting PDF page range {start_page}-{end_page}: {str(e)}")
        return f"Error extracting pages {start_page+1}-{end_page+1}: {str(e)}"

def extract_text_from_docx(file_path: str) -> str:
    """Extract text from Word documents."""
    doc = docx.Document(file_path)
    text = ""
    
    for para in doc.paragraphs:
        text += para.text + "\n"
    
    # Extract tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text += cell.text + " | "
            text += "\n"
    
    return text

def extract_text_from_xlsx(file_path: str) -> str:
    """Extract text from Excel spreadsheets."""
    workbook = openpyxl.load_workbook(file_path, data_only=True)
    text = ""
    
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        text += f"Sheet: {sheet_name}\n"
        
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
            if row_text.strip():
                text += row_text + "\n"
        
        text += "\n"
    
    return text

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PowerPoint presentations."""
    presentation = Presentation(file_path)
    text = ""
    
    for i, slide in enumerate(presentation.slides):
        text += f"Slide {i+1}:\n"
        
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
        
        text += "\n"
    
    return text

def extract_text_from_txt(file_path: str) -> str:
    """Extract text from plain text files."""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as file:
        return file.read()

async def extract_text_from_video(file_path: str) -> str:
    """
    Extract both audio and frames from video for comprehensive text extraction.
    Uses the advanced VideoProcessor for optimal results.
    """
    logger.info(f"Starting video extraction for file: {file_path}")
    
    try:
        # Extract audio from video and transcribe it
        logger.info(f"Extracting audio from video: {file_path}")
        
        # Use the existing audio transcription function
        audio_text = extract_text_from_audio(file_path)
        
        if audio_text and not audio_text.startswith("Error"):
            logger.info(f"Audio transcription completed: {len(audio_text)} characters")
            return audio_text
        else:
            logger.warning("Audio transcription failed or returned error")
            return "No audio content extracted or no audio track found"
            
    except Exception as e:
        logger.exception(f"Error processing video audio: {str(e)}")
        return f"Error extracting audio from video: {str(e)}"

def extract_text_from_video_frames(file_path: str, max_duration: float = None, interval: float = None, max_frames: int = None) -> str:
    """
    Extract frames from video at regular intervals and perform advanced OCR on them.
    Uses multiple image processing techniques to maximize text recognition accuracy.
    Applies the same successful techniques used for static image processing.
    
    Args:
        file_path: Path to the video file
        max_duration: Maximum duration in seconds to process (None for full video)
        interval: Interval between frames in seconds (overrides FRAME_INTERVAL)
        max_frames: Maximum number of frames to process (None for all frames)
    """
    logger.info(f"Starting frame extraction from: {file_path}")
    text = ""
    
    try:
        # Open the video file
        logger.debug(f"Opening video file with OpenCV: {file_path}")
        cap = cv2.VideoCapture(file_path)
        
        if not cap.isOpened():
            logger.error(f"Could not open video file: {file_path}")
            return "Error: Could not open video file."
        
        # Get video properties
        fps = max(1, int(cap.get(cv2.CAP_PROP_FPS)))  # Ensure fps is at least 1
        total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        duration = total_frames / fps if fps > 0 else 0
        logger.debug(f"Video properties for frame extraction: FPS={fps}, Total Frames={total_frames}, Duration={duration}s")
        
        # Calculate frame interval (in frames)
        # Use the specified interval or default to FRAME_INTERVAL (5 seconds)
        sampling_interval = interval if interval is not None else FRAME_INTERVAL
        # Convert seconds to number of frames
        frame_interval = int(fps * sampling_interval)
        # Ensure we always have at least 1 frame interval to prevent processing every frame
        frame_interval = max(1, frame_interval)
        
        logger.info(f"Will extract one frame every {sampling_interval} seconds (every {frame_interval} frames) at {fps} FPS")
        
        # Determine the range of frames to process
        if max_duration and fps > 0:
            max_frame_count = int(max_duration * fps)
            logger.debug(f"Limited to {max_frame_count} frames based on max_duration of {max_duration}s")
        else:
            max_frame_count = total_frames
            logger.debug(f"Processing the entire video ({total_frames} frames)")
        
        # Calculate expected number of frames to be processed
        expected_frames = max_frame_count // frame_interval + 1
        
        if max_frames:
            if expected_frames > max_frames:
                logger.debug(f"Expected to process {expected_frames} frames, limiting to {max_frames} as requested")
            else:
                logger.debug(f"Expected to process {expected_frames} frames, which is below the limit of {max_frames}")
        else:
            logger.debug(f"Expected to process approximately {expected_frames} frames (no limit specified)")
        
        # Define OCR configurations for better results
        custom_config = r'--oem 3 --psm 6'  # Default: Assume a single block of text
        alternate_config = r'--oem 3 --psm 1 -l eng'  # Automatic page segmentation with OSD
        dense_text_config = r'--oem 3 --psm 4'  # Assume a single column of text
        sparse_text_config = r'--oem 3 --psm 11'  # Sparse text with OSD
        
        # Process frames at regular intervals
        frame_count = 0
        frames_processed = 0
        frames_with_text = 0
        
        # Track timestamp for each frame's text
        frame_texts = {}
        
        # For debugging
        processed_timestamps = []
        
        logger.debug("Starting frame processing loop")
        while frame_count < max_frame_count:
            ret, frame = cap.read()
            
            if not ret:
                logger.debug(f"End of video reached after {frame_count} frames")
                break
                
            # Process frame at intervals
            if frame_count % frame_interval == 0:
                timestamp = frame_count / fps
                minutes = int(timestamp // 60)
                seconds = int(timestamp % 60)
                time_marker = f"[{minutes:02d}:{seconds:02d}]"
                
                # Store timestamp for debugging
                processed_timestamps.append(timestamp)
                
                logger.debug(f"Processing frame {frame_count} at timestamp {timestamp:.2f}s")
                
                try:
                    # Process frame using multiple techniques to maximize OCR accuracy
                    frame_ocr_results = []
                    
                    # Convert to color and apply different processing techniques
                    # These techniques mirror the successful approaches in extract_text_from_image
                    
                    # 1. Original frame (color)
                    original_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
                    pil_image = Image.fromarray(original_rgb)
                    original_text = pytesseract.image_to_string(pil_image, config=custom_config)
                    if original_text.strip():
                        frame_ocr_results.append(original_text)
                        logger.debug(f"Original image OCR: {len(original_text.strip())} chars")
                    
                    # 2. Convert to grayscale (basic)
                    gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
                    gray_pil = Image.fromarray(gray)
                    gray_text = pytesseract.image_to_string(gray_pil, config=custom_config)
                    if gray_text.strip():
                        frame_ocr_results.append(gray_text)
                        logger.debug(f"Grayscale OCR: {len(gray_text.strip())} chars")
                    
                    # 3. Apply Otsu's thresholding (very effective for text)
                    _, otsu = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY_INV | cv2.THRESH_OTSU)
                    otsu_pil = Image.fromarray(otsu)
                    otsu_text = pytesseract.image_to_string(otsu_pil, config=custom_config)
                    if otsu_text.strip():
                        frame_ocr_results.append(otsu_text)
                        logger.debug(f"Otsu threshold OCR: {len(otsu_text.strip())} chars")
                    
                    # 4. Apply regular binary thresholding
                    _, thresh = cv2.threshold(gray, 150, 255, cv2.THRESH_BINARY)
                    thresh_pil = Image.fromarray(thresh)
                    thresh_text = pytesseract.image_to_string(thresh_pil, config=custom_config)
                    if thresh_text.strip():
                        frame_ocr_results.append(thresh_text)
                        logger.debug(f"Binary threshold OCR: {len(thresh_text.strip())} chars")
                    
                    # 5. Apply adaptive thresholding for complex backgrounds
                    adaptive = cv2.adaptiveThreshold(
                        gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, 
                        cv2.THRESH_BINARY, 11, 2
                    )
                    adaptive_pil = Image.fromarray(adaptive)
                    adaptive_text = pytesseract.image_to_string(adaptive_pil, config=dense_text_config)
                    if adaptive_text.strip():
                        frame_ocr_results.append(adaptive_text)
                        logger.debug(f"Adaptive threshold OCR: {len(adaptive_text.strip())} chars")
                    
                    # 6. Apply morphological operations to enhance text
                    # First dilate to connect nearby text components
                    kernel = np.ones((1, 1), np.uint8)
                    dilated = cv2.dilate(gray, kernel, iterations=1)
                    # Then apply thresholding
                    _, dilated_thresh = cv2.threshold(dilated, 150, 255, cv2.THRESH_BINARY)
                    dilated_pil = Image.fromarray(dilated_thresh)
                    dilated_text = pytesseract.image_to_string(dilated_pil, config=dense_text_config)
                    if dilated_text.strip():
                        frame_ocr_results.append(dilated_text)
                        logger.debug(f"Dilated OCR: {len(dilated_text.strip())} chars")
                    
                    # 7. Apply distance transform for improved text detection
                    # Invert the binary image first
                    _, binary_inv = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY_INV | cv2.THRESH_OTSU)
                    # Apply distance transform
                    dist = cv2.distanceTransform(binary_inv, cv2.DIST_L2, 5)
                    # Normalize and scale to 0-255 range
                    cv2.normalize(dist, dist, 0, 255.0, cv2.NORM_MINMAX)
                    dist_img = dist.astype(np.uint8)
                    # Threshold the distance image
                    _, dist_thresh = cv2.threshold(dist_img, 50, 255, cv2.THRESH_BINARY)
                    dist_pil = Image.fromarray(dist_thresh)
                    dist_text = pytesseract.image_to_string(dist_pil, config=sparse_text_config)
                    if dist_text.strip():
                        frame_ocr_results.append(dist_text)
                        logger.debug(f"Distance transform OCR: {len(dist_text.strip())} chars")
                    
                    # 8. Apply CLAHE (Contrast Limited Adaptive Histogram Equalization)
                    clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
                    clahe_img = clahe.apply(gray)
                    # Apply thresholding on CLAHE result
                    _, clahe_thresh = cv2.threshold(clahe_img, 150, 255, cv2.THRESH_BINARY)
                    clahe_pil = Image.fromarray(clahe_thresh)
                    clahe_text = pytesseract.image_to_string(clahe_pil, config=custom_config)
                    if clahe_text.strip():
                        frame_ocr_results.append(clahe_text)
                        logger.debug(f"CLAHE OCR: {len(clahe_text.strip())} chars")
                    
                    # 9. Apply noise removal (median blur)
                    denoised = cv2.medianBlur(gray, 3)
                    # Apply thresholding on denoised image
                    _, denoised_thresh = cv2.threshold(denoised, 150, 255, cv2.THRESH_BINARY)
                    denoised_pil = Image.fromarray(denoised_thresh)
                    denoised_text = pytesseract.image_to_string(denoised_pil, config=custom_config)
                    if denoised_text.strip():
                        frame_ocr_results.append(denoised_text)
                        logger.debug(f"Denoised OCR: {len(denoised_text.strip())} chars")
                    
                    # 10. Apply edge detection and use inverted result
                    edges = cv2.Canny(gray, 100, 200)
                    edges_dilated = cv2.dilate(edges, kernel, iterations=1)
                    edges_inv = cv2.bitwise_not(edges_dilated)
                    edges_pil = Image.fromarray(edges_inv)
                    edges_text = pytesseract.image_to_string(edges_pil, config=sparse_text_config)
                    if edges_text.strip():
                        frame_ocr_results.append(edges_text)
                        logger.debug(f"Edge detection OCR: {len(edges_text.strip())} chars")
                    
                    # Only count this frame if we processed it successfully
                    frames_processed += 1
                    
                    # Combine results and remove duplicates
                    if frame_ocr_results:
                        # Sort results by length (longer texts often more complete)
                        frame_ocr_results.sort(key=len, reverse=True)
                        
                        # Remove duplicates while preserving order
                        seen = set()
                        unique_results = []
                        for result in frame_ocr_results:
                            normalized = ''.join(result.strip().lower().split())
                            if normalized and normalized not in seen and len(normalized) > 5:
                                seen.add(normalized)
                                unique_results.append(result)
                        
                        # Use the best result (typically the longest one with most content)
                        best_text = unique_results[0] if unique_results else ""
                        
                        if best_text.strip():
                            # Store text with timestamp
                            frame_texts[timestamp] = best_text.strip()
                            frames_with_text += 1
                            logger.debug(f"Frame {frame_count} contains text ({len(best_text.strip())} chars)")
                        else:
                            logger.debug(f"No meaningful text found in frame {frame_count} after multiple OCR attempts")
                    else:
                        logger.debug(f"No text found in frame {frame_count}")
                        
                except Exception as e:
                    logger.warning(f"Error processing frame {frame_count}: {str(e)}")
            
            frame_count += 1
            
            # Check if we've reached the maximum frames to process
            # Only apply max_frames limit if explicitly specified
            if max_frames and frames_processed >= max_frames:
                logger.debug(f"Reached maximum frames to process: {max_frames}")
                break
        
        # Release resources
        logger.debug("Releasing video capture resources")
        cap.release()
        
        # Log all processed timestamps for debugging
        total_frames_expected = expected_frames
        if max_frames and total_frames_expected > max_frames:
            total_frames_expected = max_frames
            
        logger.info(f"Processed {frames_processed} frames out of {total_frames_expected} expected")
        if len(processed_timestamps) > 0:
            if len(processed_timestamps) > 10:
                # Show just a few timestamps to avoid log overflow
                timestamp_sample = processed_timestamps[:3] + ["..."] + processed_timestamps[-3:]
                logger.info(f"Sample of frame timestamps (seconds): {timestamp_sample}")
            else:
                logger.info(f"Processed frames at timestamps (seconds): {processed_timestamps}")
        
        # Process results: Group nearby timestamps to avoid duplication of similar content
        # This helps with redundant text that appears in consecutive frames
        if frame_texts:
            # Sort by timestamp
            sorted_timestamps = sorted(frame_texts.keys())
            grouped_texts = []
            current_group = []
            current_group_time = None
            
            # Group texts that are within 3 seconds of each other
            time_threshold = 3.0  # seconds
            
            for ts in sorted_timestamps:
                if current_group_time is None or (ts - current_group_time) <= time_threshold:
                    current_group.append((ts, frame_texts[ts]))
                    current_group_time = ts
                else:
                    # New group
                    if current_group:
                        grouped_texts.append(current_group)
                    current_group = [(ts, frame_texts[ts])]
                    current_group_time = ts
            
            # Add the last group
            if current_group:
                grouped_texts.append(current_group)
            
            # Format text with timestamps, using the first timestamp for each group
            for group in grouped_texts:
                first_ts = group[0][0]
                minutes = int(first_ts // 60)
                seconds = int(first_ts % 60)
                
                # Get longest text from the group (usually most complete)
                group_texts = [item[1] for item in group]
                longest_text = max(group_texts, key=len) if group_texts else ""
                
                # Add to final output
                text += f"[{minutes:02d}:{seconds:02d}] {longest_text}\n\n"
        
        # Add summary
        logger.info(f"Frame processing complete: Processed {frames_processed} frames, found text in {frames_with_text} frames")
        if frames_processed > 0:
            text += f"Processed {frames_processed} frames with visual content. Found text in {frames_with_text} frames.\n"
            if len(processed_timestamps) > 0:
                total_duration = max(processed_timestamps) if processed_timestamps else 0
                text += f"Frames were extracted at {sampling_interval}-second intervals from 0 to {int(total_duration)} seconds.\n"
        else:
            logger.warning("No frames were processed from the video")
            text += "No text content could be extracted from video frames."
    
    except Exception as e:
        logger.exception(f"Error extracting frames from video: {str(e)}")
        text = f"Error extracting frames from video: {str(e)}"
    
    return text

async def extract_frames_async(file_path: str, max_duration: float = None, interval: float = None, max_frames: int = None) -> str:
    """
    Asynchronous wrapper for extracting frames from a video file.
    This allows proper timeouts and cancellation to prevent server timeouts.
    """
    return extract_text_from_video_frames(file_path, max_duration, interval, max_frames)

async def process_audio_async(audio_path: str) -> str:
    """
    Process audio file asynchronously with proper error handling.
    """
    try:
        logger.debug(f"Starting async audio processing: {audio_path}")
        return extract_text_from_audio(audio_path)
    except Exception as e:
        logger.exception(f"Error in async audio processing: {str(e)}")
        return f"Error processing audio: {str(e)}"

def extract_text_from_audio(file_path: str) -> str:
    """
    Transcribe audio to text using Whisper.
    """
    logger.info(f"Starting audio transcription for: {file_path}")
    try:
        logger.debug(f"Running Whisper transcription on: {file_path}")
        result = whisper_model.transcribe(file_path)
        logger.info(f"Audio transcription completed: {len(result['text'])} characters")
        return result["text"]
    except Exception as e:
        logger.exception(f"Error transcribing audio: {str(e)}")
        return f"Error transcribing audio: {str(e)}"

def extract_text_from_image(file_path: str) -> str:
    """
    Extract text from images using OCR with preprocessing for better results.
    Includes multiple preprocessing methods, orientation detection, and language detection.
    Uses advanced image processing techniques to maximize text extraction accuracy.
    """
    logger.info(f"Starting image OCR for: {file_path}")
    
    try:
        # Open the image
        original_image = Image.open(file_path)
        
        # Get image metadata
        width, height = original_image.size
        format_type = original_image.format
        mode = original_image.mode
        
        # Add metadata to the extracted text
        text = f"--- Image Content Extraction ---\n\n"
        text += f"Image Size: {width}x{height} pixels\n"
        text += f"Format: {format_type}\n"
        text += f"Color Mode: {mode}\n\n"
        
        # Process with multiple approaches and combine results
        ocr_results = []
        
        # Define OCR configurations for better results
        custom_config = r'--oem 3 --psm 6'  # Default: Assume a single block of text
        alternate_config = r'--oem 3 --psm 1 -l eng'  # Automatic page segmentation with OSD
        dense_text_config = r'--oem 3 --psm 4'  # Assume a single column of text
        sparse_text_config = r'--oem 3 --psm 11'  # Sparse text with OSD
        
        # Try to detect orientation and language
        try:
            logger.debug("Detecting orientation and language")
            osd_info = pytesseract.image_to_osd(original_image)
            logger.debug(f"OSD info: {osd_info}")
            
            # Extract rotation angle
            if "Rotate: " in osd_info:
                rotation_angle = int(osd_info.split("Rotate: ")[1].split("\n")[0])
                if rotation_angle != 0:
                    logger.debug(f"Detected rotation angle: {rotation_angle}")
                    rotated_image = original_image.rotate(-rotation_angle, expand=True)
                    text += f"Image Orientation: Rotated {rotation_angle} degrees (corrected)\n"
                    # Use the rotated image for further processing
                    original_image = rotated_image
        except Exception as e:
            logger.warning(f"Could not detect orientation: {str(e)}")
        
        # 1. Original image OCR
        logger.debug("Running OCR on original image")
        original_text = pytesseract.image_to_string(original_image, config=custom_config)
        if original_text.strip():
            ocr_results.append(original_text)
            logger.debug(f"Original image OCR: {len(original_text.strip())} chars")
        
        # Try alternative PSM for original image
        alt_original_text = pytesseract.image_to_string(original_image, config=alternate_config)
        if alt_original_text.strip() and alt_original_text != original_text:
            ocr_results.append(alt_original_text)
            logger.debug(f"Original image with alt config: {len(alt_original_text.strip())} chars")
        
        # 2. Convert to OpenCV format for preprocessing
        if mode == "RGB":
            opencv_image = cv2.cvtColor(np.array(original_image), cv2.COLOR_RGB2BGR)
        else:
            # Convert to RGB first if not already
            rgb_image = original_image.convert('RGB')
            opencv_image = cv2.cvtColor(np.array(rgb_image), cv2.COLOR_RGB2BGR)
        
        # 3. Grayscale conversion (basic)
        logger.debug("Running OCR on grayscale image")
        gray_image = cv2.cvtColor(opencv_image, cv2.COLOR_BGR2GRAY)
        gray_pil = Image.fromarray(gray_image)
        gray_text = pytesseract.image_to_string(gray_pil, config=custom_config)
        if gray_text.strip():
            ocr_results.append(gray_text)
            logger.debug(f"Grayscale OCR: {len(gray_text.strip())} chars")
        
        # 4. Apply Otsu's thresholding (very effective for text)
        logger.debug("Running OCR on Otsu thresholded image")
        _, otsu = cv2.threshold(gray_image, 0, 255, cv2.THRESH_BINARY_INV | cv2.THRESH_OTSU)
        otsu_pil = Image.fromarray(otsu)
        otsu_text = pytesseract.image_to_string(otsu_pil, config=custom_config)
        if otsu_text.strip():
            ocr_results.append(otsu_text)
            logger.debug(f"Otsu threshold OCR: {len(otsu_text.strip())} chars")
        
        # 5. Apply regular binary thresholding
        logger.debug("Running OCR on binary thresholded image")
        _, threshold_img = cv2.threshold(gray_image, 150, 255, cv2.THRESH_BINARY)
        threshold_pil = Image.fromarray(threshold_img)
        threshold_text = pytesseract.image_to_string(threshold_pil, config=custom_config)
        if threshold_text.strip():
            ocr_results.append(threshold_text)
            logger.debug(f"Binary threshold OCR: {len(threshold_text.strip())} chars")
        
        # 6. Apply adaptive thresholding for complex backgrounds
        logger.debug("Running OCR on adaptive thresholded image")
        adaptive_threshold = cv2.adaptiveThreshold(
            gray_image, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 11, 2
        )
        adaptive_pil = Image.fromarray(adaptive_threshold)
        adaptive_text = pytesseract.image_to_string(adaptive_pil, config=dense_text_config)
        if adaptive_text.strip():
            ocr_results.append(adaptive_text)
            logger.debug(f"Adaptive threshold OCR: {len(adaptive_text.strip())} chars")
        
        # 7. Apply distance transform for improved text detection
        logger.debug("Running OCR with distance transform")
        # Invert the binary image first
        _, binary_inv = cv2.threshold(gray_image, 0, 255, cv2.THRESH_BINARY_INV | cv2.THRESH_OTSU)
        # Apply distance transform
        dist = cv2.distanceTransform(binary_inv, cv2.DIST_L2, 5)
        # Normalize and scale to 0-255 range
        cv2.normalize(dist, dist, 0, 255.0, cv2.NORM_MINMAX)
        dist_img = dist.astype(np.uint8)
        # Threshold the distance image
        _, dist_thresh = cv2.threshold(dist_img, 50, 255, cv2.THRESH_BINARY)
        dist_pil = Image.fromarray(dist_thresh)
        dist_text = pytesseract.image_to_string(dist_pil, config=sparse_text_config)
        if dist_text.strip():
            ocr_results.append(dist_text)
            logger.debug(f"Distance transform OCR: {len(dist_text.strip())} chars")
            
        # 8. Apply morphological operations to enhance text
        logger.debug("Running OCR with morphological operations")
        # Create kernel for morphological operations
        kernel = np.ones((1, 1), np.uint8)
        # First dilate to connect nearby text components
        dilated = cv2.dilate(gray_image, kernel, iterations=1)
        # Then apply thresholding
        _, dilated_thresh = cv2.threshold(dilated, 150, 255, cv2.THRESH_BINARY)
        dilated_pil = Image.fromarray(dilated_thresh)
        dilated_text = pytesseract.image_to_string(dilated_pil, config=dense_text_config)
        if dilated_text.strip():
            ocr_results.append(dilated_text)
            logger.debug(f"Dilated OCR: {len(dilated_text.strip())} chars")
            
        # 9. Apply noise removal for better text detection
        logger.debug("Running OCR with noise removal")
        denoised = cv2.medianBlur(gray_image, 3)
        # Apply thresholding on denoised image
        _, denoised_thresh = cv2.threshold(denoised, 150, 255, cv2.THRESH_BINARY)
        denoised_pil = Image.fromarray(denoised_thresh)
        denoised_text = pytesseract.image_to_string(denoised_pil, config=custom_config)
        if denoised_text.strip():
            ocr_results.append(denoised_text)
            logger.debug(f"Denoised OCR: {len(denoised_text.strip())} chars")
            
        # 10. Apply edge enhancement to detect text boundaries better
        logger.debug("Enhancing edges for text boundary detection")
        edges = cv2.Canny(gray_image, 100, 200)
        edges_dilated = cv2.dilate(edges, kernel, iterations=1)
        # Invert the image for better OCR (black text on white background)
        edges_inv = cv2.bitwise_not(edges_dilated)
        edges_pil = Image.fromarray(edges_inv)
        edges_text = pytesseract.image_to_string(edges_pil, config=sparse_text_config)
        if edges_text.strip():
            ocr_results.append(edges_text)
            logger.debug(f"Edge detection OCR: {len(edges_text.strip())} chars")
        
        # 11. Apply CLAHE for better contrast
        logger.debug("Enhancing image contrast with CLAHE")
        clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
        clahe_img = clahe.apply(gray_image)
        # Apply thresholding on CLAHE result
        _, clahe_thresh = cv2.threshold(clahe_img, 150, 255, cv2.THRESH_BINARY)
        clahe_pil = Image.fromarray(clahe_thresh)
        clahe_text = pytesseract.image_to_string(clahe_pil, config=custom_config)
        if clahe_text.strip():
            ocr_results.append(clahe_text)
            logger.debug(f"CLAHE OCR: {len(clahe_text.strip())} chars")
        
        # 12. For high-res images, try scaled down version
        if width > 2000 or height > 2000:
            logger.debug("Image is high resolution, trying scaled down version")
            scale_percent = 50  # percent of original size
            new_width = int(width * scale_percent / 100)
            new_height = int(height * scale_percent / 100)
            resized_img = cv2.resize(gray_image, (new_width, new_height))
            resized_pil = Image.fromarray(resized_img)
            resized_text = pytesseract.image_to_string(resized_pil, config=custom_config)
            if resized_text.strip():
                ocr_results.append(resized_text)
                logger.debug(f"Resized OCR: {len(resized_text.strip())} chars")
        
        # Combine results, removing duplicates
        # Sort results by length (longer texts are often more complete)
        ocr_results.sort(key=len, reverse=True)
        
        # Remove duplicates while preserving order
        seen = set()
        unique_results = []
        for result in ocr_results:
            # Normalize text for comparison (remove whitespace and lowercase)
            normalized = ''.join(result.strip().lower().split())
            if normalized and normalized not in seen and len(normalized) > 5:  # Ensure meaningful content
                seen.add(normalized)
                unique_results.append(result)
        
        combined_text = "\n\n".join(unique_results[:3])  # Keep top 3 most complete results
        
        if combined_text.strip():
            text += "--- OCR Text Content ---\n\n"
            text += combined_text
            logger.info(f"Successfully extracted {len(combined_text.strip())} chars from image")
        else:
            text += "No text was detected in this image."
            logger.warning("No text detected in the image after multiple OCR attempts")
        
        return text
        
    except Exception as e:
        logger.exception(f"Error extracting text from image: {str(e)}")
        return f"Error extracting text from image: {str(e)}"

def create_chunks(text: str, document_id: str, repository_id: str, user_id: str) -> List[LangchainDocument]:
    """
    Split text into chunks for processing and embedding.
    Uses a smaller chunk size for dense texts to ensure we stay within token limits.
    
    Args:
        text: The text to split into chunks
        document_id: ID of the document
        repository_id: ID of the repository
        user_id: ID of the user who owns the document
        
    Returns:
        List of LangchainDocument chunks
    """
    # Estimate token density (tokens per character)
    token_estimate = len(tokenizer.encode(text[:10000] if len(text) > 10000 else text))
    char_count = min(len(text), 10000)
    token_density = token_estimate / char_count if char_count > 0 else 0.25
    
    logger.info(f"Estimated token density: {token_density:.4f} tokens per character")
    
    # Adjust chunk size based on token density to avoid excessive tokens per chunk
    # Higher density means we need smaller chunks
    adjusted_chunk_size = min(CHUNK_SIZE, int(1500 / token_density)) if token_density > 0 else CHUNK_SIZE
    adjusted_chunk_size = max(adjusted_chunk_size, 300)  # Ensure chunk size isn't too small
    
    logger.info(f"Using adjusted chunk size of {adjusted_chunk_size} characters (standard: {CHUNK_SIZE})")
    
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=adjusted_chunk_size,
        chunk_overlap=min(CHUNK_OVERLAP, adjusted_chunk_size // 5),  # Adjust overlap proportionally
        length_function=len,
    )
    
    chunks = text_splitter.create_documents([text])
    
    # Add document_id, repository_id, and user_id to metadata
    for i, chunk in enumerate(chunks):
        chunk.metadata["document_id"] = document_id
        chunk.metadata["repository_id"] = repository_id
        chunk.metadata["user_id"] = user_id
        chunk.metadata["chunk_id"] = f"{document_id}_{i}"
    
    logger.info(f"Created {len(chunks)} chunks from document")
    
    return chunks

def count_tokens_in_chunks(chunks: List[LangchainDocument]) -> int:
    """
    Count the approximate number of tokens in a list of chunks.
    
    Args:
        chunks: List of LangchainDocument chunks
        
    Returns:
        Approximate token count
    """
    total_tokens = 0
    sample_size = min(20, len(chunks))  # Sample up to 20 chunks for efficiency
    
    if sample_size > 0:
        # Calculate average tokens per chunk from the sample
        sampled_chunks = chunks[:sample_size]
        sampled_tokens = sum(len(tokenizer.encode(chunk.page_content)) for chunk in sampled_chunks)
        avg_tokens_per_chunk = sampled_tokens / sample_size
        
        # Estimate total tokens
        total_tokens = int(avg_tokens_per_chunk * len(chunks))
    
    return total_tokens

async def process_chunks_in_batches(chunks: List[LangchainDocument], document_id: str):
    """
    Process chunks in batches to avoid exceeding OpenAI's token limits.
    
    Args:
        chunks: List of LangchainDocument chunks
        document_id: ID of the document
    """
    total_chunks = len(chunks)
    logger.info(f"Processing {total_chunks} chunks in batches for document {document_id}")
    
    # Process in batches to avoid token limits
    batch_size = MAX_CHUNKS_PER_BATCH  # Process a maximum of 100 chunks per batch
    total_batches = (total_chunks + batch_size - 1) // batch_size
    
    for batch_idx in range(total_batches):
        start_idx = batch_idx * batch_size
        end_idx = min(start_idx + batch_size, total_chunks)
        current_batch = chunks[start_idx:end_idx]
        
        # Estimate tokens in current batch
        batch_tokens = sum(len(tokenizer.encode(chunk.page_content)) for chunk in current_batch)
        
        logger.info(f"Processing batch {batch_idx+1}/{total_batches} with {len(current_batch)} chunks ({batch_tokens} tokens)")
        
        # If the batch might exceed token limits, process in smaller sub-batches
        if batch_tokens > MAX_TOKENS_PER_REQUEST:
            logger.warning(f"Batch {batch_idx+1} exceeds token limit, splitting into smaller batches")
            await process_oversized_batch(current_batch, document_id, batch_idx, total_batches)
        else:
            # Process the batch normally
            try:
                await add_chunks_to_vectorstore(current_batch, document_id)
                logger.info(f"Successfully processed batch {batch_idx+1}/{total_batches}")
            except Exception as e:
                logger.error(f"Error processing batch {batch_idx+1}: {str(e)}")
                # Try to process this batch in smaller chunks
                await process_oversized_batch(current_batch, document_id, batch_idx, total_batches)
        
        # Small delay between batches to avoid rate limiting
        await asyncio.sleep(0.5)
    
    logger.info(f"Completed processing all {total_chunks} chunks for document {document_id}")

async def process_oversized_batch(chunks: List[LangchainDocument], document_id: str, batch_idx: int, total_batches: int):
    """
    Process an oversized batch by breaking it into smaller sub-batches.
    
    Args:
        chunks: List of chunks that collectively exceed the token limit
        document_id: ID of the document
        batch_idx: Current batch index for logging
        total_batches: Total number of main batches for logging
    """
    current_batch = []
    current_tokens = 0
    sub_batch_idx = 1
    
    for chunk in chunks:
        chunk_tokens = len(tokenizer.encode(chunk.page_content))
        
        # If adding this chunk would exceed the token limit, process the current batch
        if current_tokens + chunk_tokens > MAX_TOKENS_PER_REQUEST or len(current_batch) >= 20:
            if current_batch:
                logger.info(f"Processing sub-batch {sub_batch_idx} of oversized batch {batch_idx+1}/{total_batches} ({current_tokens} tokens)")
                try:
                    await add_chunks_to_vectorstore(current_batch, document_id)
                    logger.info(f"Successfully processed sub-batch {sub_batch_idx}")
                except Exception as e:
                    logger.error(f"Error processing sub-batch {sub_batch_idx}: {str(e)}")
                    # If a small batch still fails, process chunks individually
                    await process_individual_chunks(current_batch, document_id)
                
                # Reset for next sub-batch
                current_batch = []
                current_tokens = 0
                sub_batch_idx += 1
                
                # Small delay between sub-batches
                await asyncio.sleep(0.5)
        
        # Add chunk to current batch
        current_batch.append(chunk)
        current_tokens += chunk_tokens
    
    # Process any remaining chunks
    if current_batch:
        logger.info(f"Processing final sub-batch {sub_batch_idx} of oversized batch {batch_idx+1}/{total_batches} ({current_tokens} tokens)")
        try:
            await add_chunks_to_vectorstore(current_batch, document_id)
            logger.info(f"Successfully processed final sub-batch")
        except Exception as e:
            logger.error(f"Error processing final sub-batch: {str(e)}")
            # If a small batch still fails, process chunks individually
            await process_individual_chunks(current_batch, document_id)

async def process_individual_chunks(chunks: List[LangchainDocument], document_id: str):
    """
    Process chunks individually as a last resort.
    
    Args:
        chunks: List of chunks to process individually
        document_id: ID of the document
    """
    logger.warning(f"Processing {len(chunks)} chunks individually")
    
    success_count = 0
    failure_count = 0
    
    for i, chunk in enumerate(chunks):
        try:
            # Process one chunk at a time
            await add_chunks_to_vectorstore([chunk], document_id)
            success_count += 1
            
            # Add a small delay between requests
            await asyncio.sleep(0.2)
        except Exception as e:
            failure_count += 1
            logger.error(f"Failed to process individual chunk {i}: {str(e)}")
            
            # If too many failures, stop processing
            if failure_count > 5:
                logger.error("Too many failures when processing individual chunks, stopping")
                break
    
    logger.info(f"Individual chunk processing complete: {success_count} succeeded, {failure_count} failed") 